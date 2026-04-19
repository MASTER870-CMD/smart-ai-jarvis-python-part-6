import os
import re
from pptx import Presentation
from pptx.util import Pt

def generate_ppt(topic, client, model_id):
    """
    Generates a comprehensive, content-rich PowerPoint with 7 slides.
    Requires a 'Templates' folder with 'Professional.pptx' inside it.
    """
    try:
        # ====================================================
        # 1. AI "Architect" Step: Plan the content & Choose Theme
        # ====================================================
        print(f"DEBUG: Generating PPT content for {topic}...")
        
        sys_prompt = """
        You are a senior corporate strategist creating a high-stakes presentation. 
        Your goal is to provide DEEP, DETAILED insights. 
        
        Rules:
        1. Choose a theme: 'Professional' (business) or 'Creative' (tech/art).
        2. Generate 7 SLIDES.
        3. For bullet points, write FULL SENTENCES (15-20 words each) explaining the "Why" and "How". Do not use short keywords.
        4. Make the content dense and informative.
        
        Output format strictly like this:
        THEME: Professional (or Creative)
        SLIDE 1: Title | Subtitle
        SLIDE 2: Title | Detailed Point 1; Detailed Point 2; Detailed Point 3; Detailed Point 4
        SLIDE 3: Title | Detailed Point 1; Detailed Point 2; Detailed Point 3
        ...
        SLIDE 7: Conclusion | Final summary paragraph
        """
        
        completion = client.chat.completions.create(
            model=model_id,
            messages=[
                {"role": "system", "content": sys_prompt},
                {"role": "user", "content": f"Create a detailed presentation about: {topic}"}
            ]
        )
        response_text = completion.choices[0].message.content
        
        # ====================================================
        # 2. Parse the AI's response
        # ====================================================
        lines = response_text.split('\n')
        theme_choice = "Professional" # Default fallback
        slides_content = []
        
        for line in lines:
            line = line.strip()
            if line.startswith("THEME:"):
                theme_choice = line.replace("THEME:", "").strip()
            elif line.startswith("SLIDE"):
                # Remove "SLIDE X:" and split by "|"
                clean_line = re.sub(r'SLIDE \d+:', '', line).strip()
                if "|" in clean_line:
                    parts = clean_line.split("|")
                    title = parts[0].strip()
                    content = parts[1].strip()
                    slides_content.append((title, content))

        # ====================================================
        # 3. Load the Smart Template (With Error Checking)
        # ====================================================
        # This expects a folder named 'Templates' next to this script
        template_path = os.path.join("Templates", f"{theme_choice}.pptx")
        
        print(f"DEBUG: Looking for template at: {os.path.abspath(template_path)}")
        
        if os.path.exists(template_path):
            prs = Presentation(template_path)
            print(f"DEBUG: Successfully loaded {theme_choice} theme.")
        else:
            # IMPORTANT: If file is missing, we stop and warn the user
            return f"ERROR: Could not find '{theme_choice}.pptx' in the Templates folder. Please check the file name."

        # ====================================================
        # 4. Build the Slides
        # ====================================================
        for i, (title_text, content_text) in enumerate(slides_content):
            # Slide 1 (Index 0) is usually Title Layout
            # Slide 2+ (Index 1) is usually Content Layout
            layout_index = 0 if i == 0 else 1 
            
            # Safety check if template has fewer layouts
            if layout_index >= len(prs.slide_layouts): layout_index = 0
            
            slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
            
            # --- Set Title ---
            if slide.shapes.title:
                slide.shapes.title.text = title_text
            
            # --- Set Content ---
            # We look for the 'Body' placeholder (usually index 1)
            body_shape = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    body_shape = shape
                    break
            
            if body_shape:
                tf = body_shape.text_frame
                tf.text = "" # Clear default dummy text
                
                # Split content by ";" to make bullet points
                bullets = content_text.split(";")
                for bullet in bullets:
                    p = tf.add_paragraph()
                    p.text = bullet.strip()
                    p.level = 0
                    
                    # Formatting for readability
                    p.space_after = Pt(14)  # Space between bullets
                    p.font.size = Pt(18)    # Readable font size

        # ====================================================
        # 5. Save & Open
        # ====================================================
        desktop = os.path.join(os.path.expanduser("~"), "Desktop")
        output_folder = os.path.join(desktop, "Jarvis_PPT_Generated")
        
        if not os.path.exists(output_folder):
            os.makedirs(output_folder)
            
        # Clean up filename (remove spaces/special chars)
        safe_topic = "".join([c for c in topic if c.isalnum() or c in (' ', '_')]).strip().replace(" ", "_")
        filename = f"{safe_topic}_{theme_choice}_Detailed.pptx"
        save_path = os.path.join(output_folder, filename)
        
        prs.save(save_path)
        print(f"DEBUG: Saved to {save_path}")
        
        # Automatically open the file
        os.startfile(save_path) 
        
        return f"Presentation generated successfully: {filename}"

    except Exception as e:
        print(f"CRITICAL ERROR: {str(e)}")
        return f"Failed to generate PPT: {str(e)}"